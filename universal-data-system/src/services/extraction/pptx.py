"""Content extraction for PPTX files."""

import io
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from src.utils.logging import get_logger

logger = get_logger(__name__)


class PptxExtractor:
    """Extract text content from PPTX files."""

    def extract(self, content: bytes) -> str:
        """
        Extract text from PPTX bytes.

        Returns Markdown-formatted text with slide structure.
        """
        try:
            prs = Presentation(io.BytesIO(content))

            slides = []
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_parts = [f"## Slide {slide_num}"]

                for shape in slide.shapes:
                    text = self._extract_shape_text(shape)
                    if text:
                        slide_parts.append(text)

                # Extract notes if present
                if slide.has_notes_slide:
                    notes = slide.notes_slide.notes_text_frame.text.strip()
                    if notes:
                        slide_parts.append(f"\n> **Notes:** {notes}")

                slides.append("\n\n".join(slide_parts))

            result = "\n\n---\n\n".join(slides)
            logger.debug("Extracted PPTX", slides=len(prs.slides), chars=len(result))
            return result

        except Exception as e:
            logger.error("PPTX extraction failed", error=str(e))
            raise

    def _extract_shape_text(self, shape) -> str:
        """Extract text from a shape."""
        parts = []

        # Check if shape has text frame
        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                text = para.text.strip()
                if text:
                    # Check if it looks like a title
                    if para.level == 0 and len(text) < 100:
                        parts.append(f"### {text}")
                    else:
                        parts.append(text)

        # Handle tables
        if shape.has_table:
            parts.append(self._table_to_markdown(shape.table))

        return "\n".join(parts)

    def _table_to_markdown(self, table) -> str:
        """Convert a PPTX table to Markdown format."""
        rows = []

        for i, row in enumerate(table.rows):
            cells = []
            for cell in row.cells:
                text = cell.text.strip().replace("|", "\\|")
                cells.append(text)
            rows.append("| " + " | ".join(cells) + " |")

            # Add header separator after first row
            if i == 0:
                separator = "| " + " | ".join(["---"] * len(cells)) + " |"
                rows.append(separator)

        return "\n".join(rows)
